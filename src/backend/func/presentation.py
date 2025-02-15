import streamlit as st
from pptx import Presentation
from pptx.util import Inches

# Function to generate PPT
def create_ppt(pitch_text):
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]  # Title Slide Layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Startup Pitch"
    subtitle.text = "Generated Presentation"

    # Process the user input into slides
    key_points = pitch_text.split("\n")  # Splitting into bullet points
    
    for point in key_points:
        if point.strip():  # Avoid empty lines
            slide_layout = prs.slide_layouts[1]  # Title & Content Layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = "Key Point"
            content.text = point

    # Save presentation
    ppt_file = "generated_presentation.pptx"
    prs.save(ppt_file)
    return ppt_file
    
# Streamlit UI
st.title("AI-Powered Pitch Deck Generator")

# User Input
prompt = st.text_area("Enter the key points for your pitch deck:")

if st.button("Generate Presentation"):
    if prompt.strip():
        ppt_path = create_ppt(prompt)
        with open(ppt_path, "rb") as f:
            st.download_button("Download Presentation", f, file_name="Pitch_Deck.pptx")
    else:
        st.warning("Please enter some content before generating the PPT.")
